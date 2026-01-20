"""File upload component"""
import streamlit as st
from core.storage_manager import CloudStorageManager
from config.settings import SUPPORTED_FILE_TYPES
from config.languages import get_text




import os
os.environ.setdefault("PDX_EAGER_INITIALIZATION", "0")

import importlib

import streamlit as st
import requests
import json
import os
import tempfile
from PIL import Image
import pandas as pd
import matplotlib.pyplot as plt
import io
import csv
from docx import Document  # 处理Word
from PyPDF2 import PdfReader  # 处理PDF
from openpyxl import load_workbook  # 处理Excel
from pptx import Presentation  # 处理PPT
from paddleocr import PaddleOCR  # 替换pytesseract为PaddleOCR

# 初始化PaddleOCR（支持中英文）
ocr = PaddleOCR(use_angle_cls=True, lang="ch")

# 设置页面配置
st.set_page_config(page_title="AI文件处理助手", layout="wide")

# 确保中文显示正常
plt.rcParams["font.family"] = ["SimHei", "WenQuanYi Micro Hei", "Heiti TC"]

# 初始化会话状态
if 'uploaded_files' not in st.session_state:
    st.session_state.uploaded_files = {}  # 存储上传的文件信息

if 'active_file' not in st.session_state:
    st.session_state.active_file = None  # 当前选中的文件

# Deepseek API配置
DEEPSEEK_API_KEY = "sk-428249a366a7472f9403396d3c298f10"
DEEPSEEK_API_URL = "https://api.deepseek.com/v1/chat/completions"


def call_deepseek_api(prompt, system_message=None):
    """调用Deepseek API获取响应"""
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {DEEPSEEK_API_KEY}"
    }

    messages = []
    if system_message:
        messages.append({"role": "system", "content": system_message})
    messages.append({"role": "user", "content": prompt})

    data = {
        "model": "deepseek-chat",
        "messages": messages,
        "temperature": 0.7,
        "max_tokens": 1024
    }

    try:
        response = requests.post(DEEPSEEK_API_URL, headers=headers, json=data)
        response.raise_for_status()
        return response.json()['choices'][0]['message']['content']
    except Exception as e:
        st.error(f"API调用错误: {str(e)}")
        return None


def paddle_ocr_processing(image_path):
    """使用PaddleOCR进行文本提取"""
    try:
        # 调用PaddleOCR进行识别（移除cls=True参数）
        result = ocr.ocr(image_path)  # 这里删除了cls=True
        # 提取文本内容
        text = ""
        # 兼容不同版本的PaddleOCR返回格式
        if result is not None:
            for line in result:
                if line is not None:  # 处理可能的空行
                    for word_info in line:
                        if len(word_info) >= 2 and isinstance(word_info[1], tuple):
                            text += word_info[1][0] + "\n"
        return text
    except Exception as e:
        st.error(f"PaddleOCR处理错误: {str(e)}")
        return None


def save_uploaded_file(uploaded_file):
    """保存上传的文件并返回文件信息"""
    # 创建临时文件
    temp_dir = tempfile.TemporaryDirectory()
    file_path = os.path.join(temp_dir.name, uploaded_file.name)

    with open(file_path, "wb") as f:
        f.write(uploaded_file.getbuffer())

    # 获取文件扩展名
    file_ext = os.path.splitext(uploaded_file.name)[1].lower()

    # 初始化文件信息
    file_info = {
        "name": uploaded_file.name,
        "type": uploaded_file.type,
        "size": uploaded_file.size,
        "path": file_path,
        "temp_dir": temp_dir,
        "ext": file_ext,  # 新增：文件扩展名
        "labels": [],
        "extracted_text": "",  # 统一使用extracted_text存储提取的文本
        "abstract": "",
        "report": "",
        "visualization": None
    }

    return file_info


def extract_text_from_file(file_info):
    """根据文件格式提取文本内容"""
    file_path = file_info["path"]
    file_ext = file_info["ext"]
    extracted_text = ""

    try:
        # 文本文件
        if file_ext in ['.txt']:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                extracted_text = f.read()

        # CSV文件
        elif file_ext in ['.csv']:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                reader = csv.reader(f)
                for row in reader:
                    extracted_text += ", ".join(row) + "\n"

        # PDF文件
        elif file_ext in ['.pdf']:
            reader = PdfReader(file_path)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    extracted_text += page_text + "\n"

        # Word文件
        elif file_ext in ['.docx']:
            doc = Document(file_path)
            for para in doc.paragraphs:
                extracted_text += para.text + "\n"

        # Excel文件
        elif file_ext in ['.xlsx', '.xls']:
            wb = load_workbook(file_path, read_only=True)
            for sheet in wb.sheetnames:
                extracted_text += f"=== 工作表: {sheet} ===\n"
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    extracted_text += ", ".join(row_text) + "\n"
            wb.close()

        # PPT文件
        elif file_ext in ['.pptx']:
            prs = Presentation(file_path)
            for slide_idx, slide in enumerate(prs.slides, 1):
                extracted_text += f"=== 幻灯片 {slide_idx} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"

        # 图像文件（使用OCR）
        elif file_ext in ['.png', '.jpg', '.jpeg', '.bmp', '.gif']:
            extracted_text = paddle_ocr_processing(file_path)

        # 其他无法直接转换的格式（尝试OCR）
        else:
            st.warning(f"不支持的文件格式: {file_ext}，尝试使用OCR提取文本...")
            try:
                # 尝试以图像方式打开
                img = Image.open(file_path)
                extracted_text = paddle_ocr_processing(file_path)
            except:
                extracted_text = f"无法提取文本 (文件格式: {file_ext})"

    except Exception as e:
        st.error(f"提取文本时出错: {str(e)}")
        extracted_text = f"文本提取失败: {str(e)}"

    return extracted_text




def render_upload_section(storage_manager: CloudStorageManager):
    
    # 主界面
    st.title("AI 文件智能处理助手")
    st.write("上传文件并使用AI功能进行处理、分析和可视化")
    
    # 文件上传区域
    uploaded_file = st.file_uploader("点击上传文件", type=None)
    if uploaded_file is not None:
        if uploaded_file.name not in st.session_state.uploaded_files:
            # 保存文件并初始化信息
            file_info = save_uploaded_file(uploaded_file)
            st.session_state.uploaded_files[uploaded_file.name] = file_info
            st.session_state.active_file = uploaded_file.name
    
            # 自动提取文本
            with st.spinner(f"正在提取 '{uploaded_file.name}' 的文本内容..."):
                text = extract_text_from_file(file_info)
                st.session_state.uploaded_files[uploaded_file.name]["extracted_text"] = text
    
            st.success(f"文件 '{uploaded_file.name}' 上传成功!")
    
    # 显示已上传的文件
    if st.session_state.uploaded_files:
        st.subheader("已上传文件")
    
        # 显示所有标签供筛选
        all_labels = set()
        for file_info in st.session_state.uploaded_files.values():
            all_labels.update(file_info["labels"])
    
        if all_labels:
            col_filter, _ = st.columns([1, 3])
            with col_filter:
                selected_label = st.selectbox("按标签筛选文件", ["所有文件"] + list(all_labels))
    
                # 标签筛选逻辑
                filtered_files = []
                if selected_label == "所有文件":
                    filtered_files = list(st.session_state.uploaded_files.keys())
                else:
                    for file_name, file_info in st.session_state.uploaded_files.items():
                        if selected_label in file_info["labels"]:
                            filtered_files.append(file_name)
        else:
            filtered_files = list(st.session_state.uploaded_files.keys())
    
        # 选择要处理的文件
        if filtered_files:
            selected_file = st.selectbox("选择要处理的文件", filtered_files)
            st.session_state.active_file = selected_file
            file_info = st.session_state.uploaded_files[selected_file]
    
            # 显示文件基本信息
            st.write(f"### 当前处理: {file_info['name']}")
            col1, col2 = st.columns(2)
    
            with col1:
                st.write("**文件信息**")
                st.write(f"类型: {file_info['type']}")
                st.write(f"格式: {file_info['ext']}")
                st.write(f"大小: {file_info['size']} bytes")
    
                # 显示标签
                if file_info["labels"]:
                    st.write("**标签:**")
                    label_cols = st.columns(len(file_info["labels"]))
                    for i, label in enumerate(file_info["labels"]):
                        with label_cols[i]:
                            if st.button(label, key=f"label_{file_info['name']}_{label}"):
                                # 点击标签筛选文件
                                st.session_state.selected_label = label
                                st.rerun()
    
                # 显示摘要
                if file_info["abstract"]:
                    st.write("**摘要:**")
                    st.write(file_info["abstract"])
    
                st.markdown("---")
                st.write("### AI 功能")
    
                # 功能1: 自动文件识别与分类
                if st.button("🔍 自动识别与分类", key="btn_classify"):
                    with st.spinner("正在进行文件分类..."):
                        # 使用提取的文本
                        file_content = file_info["extracted_text"]
    
                        if not file_content:
                            st.error("无法获取文件内容，无法进行分类")
                        else:
                            # 调用LLM进行分类
                            prompt = f"""请分析以下文件内容，为其分配最合适的3个分类标签。
                            标签应简洁明了，用逗号分隔。只需返回标签，不要其他内容。
    
                            文件内容: {file_content[:1000]}"""
    
                            labels = call_deepseek_api(prompt)
    
                            if labels:
                                # 处理标签
                                new_labels = [label.strip() for label in labels.split(',') if label.strip()]
                                # 去重
                                unique_labels = list(set(new_labels))
                                # 添加到文件信息
                                st.session_state.uploaded_files[selected_file]["labels"].extend(
                                    [l for l in unique_labels if l not in file_info["labels"]]
                                )
                                st.success(f"已添加标签: {', '.join(unique_labels)}")
                                st.rerun()
    
                # 功能2: 文本编辑与NLP关键信息提取
                if st.button("📝 文本编辑与摘要提取", key="btn_ocr"):
                    with st.spinner("正在进行文本处理和信息提取..."):
                        # 获取已提取的文本
                        content = file_info["extracted_text"]
    
                        if not content:
                            st.error("无法获取文件内容，无法生成摘要")
                        else:
                            # 生成摘要
                            prompt = f"""请分析以下内容，生成一个简洁的摘要（约150字），
                            突出关键信息和核心内容。
    
                            内容: {content[:2000]}"""
    
                            abstract = call_deepseek_api(prompt)
                            if abstract:
                                st.session_state.uploaded_files[selected_file]["abstract"] = abstract
                                st.success("摘要生成成功!")
                                st.rerun()
    
                # 功能3: AI生成简化报告和可视化
                if st.button("📊 生成报告与可视化", key="btn_report"):
                    with st.spinner("正在生成报告和可视化..."):
                        # 获取文件内容
                        content = file_info["extracted_text"]
    
                        if not content:
                            st.error("无法获取文件内容，无法生成报告")
                        else:
                            # 生成报告
                            prompt_report = f"""请分析以下内容，生成一份简化的分析报告（约300字），
                            包括主要内容、关键发现和重要结论。
    
                            内容: {content[:3000]}"""
    
                            report = call_deepseek_api(prompt_report)
    
                            # 生成可视化数据
                            prompt_vis = f"""请分析以下内容，提取可以可视化的数据，
                            以JSON格式返回，例如：{{"类别": ["A", "B", "C"], "数值": [10, 20, 30]}}
                            确保JSON格式正确，不要包含其他内容。
    
                            内容: {content[:2000]}"""
    
                            vis_data = call_deepseek_api(prompt_vis)
    
                            # 保存结果
                            if report:
                                st.session_state.uploaded_files[selected_file]["report"] = report
    
                            if vis_data:
                                try:
                                    vis_json = json.loads(vis_data)
                                    st.session_state.uploaded_files[selected_file]["visualization"] = vis_json
                                except:
                                    st.warning("无法解析可视化数据，可能格式不正确")
    
                            st.success("报告和可视化生成成功!")
                            st.rerun()
    
            with col2:
                # 显示提取的文本（可编辑）
                if file_info["extracted_text"]:
                    st.write("**提取的文本内容（可编辑）：**")
                    edited_text = st.text_area(
                        "",
                        file_info["extracted_text"],
                        height=300,
                        key="text_editor"
                    )
                    if edited_text != file_info["extracted_text"]:
                        st.session_state.uploaded_files[selected_file]["extracted_text"] = edited_text
                        st.info("文本已更新")
    
                # 显示报告
                if file_info["report"]:
                    st.markdown("---")
                    st.write("**AI生成报告：**")
                    st.write(file_info["report"])
    
                # 显示可视化
                if file_info["visualization"]:
                    st.markdown("---")
                    st.write("**数据可视化：**")
                    try:
                        vis_data = file_info["visualization"]
    
                        # 尝试创建图表
                        fig, ax = plt.subplots(figsize=(10, 6))
    
                        # 支持多种图表类型
                        if "类别" in vis_data and "数值" in vis_data:
                            ax.bar(vis_data["类别"], vis_data["数值"])
                            ax.set_title("数据分布")
                            ax.set_xlabel("类别")
                            ax.set_ylabel("数值")
                        elif "x" in vis_data and "y" in vis_data:
                            ax.plot(vis_data["x"], vis_data["y"], marker='o')
                            ax.set_title("趋势分析")
                            ax.set_xlabel("X轴")
                            ax.set_ylabel("Y轴")
                        elif isinstance(vis_data, dict) and len(vis_data) > 0:
                            ax.pie(vis_data.values(), labels=vis_data.keys(), autopct='%1.1f%%')
                            ax.set_title("占比分析")
    
                        st.pyplot(fig)
                    except Exception as e:
                        st.error(f"可视化失败: {str(e)}")
                        st.write("原始数据:", file_info["visualization"])
    else:
        st.info("请上传文件开始处理")
    
    # 页脚信息
    st.markdown("---")
    st.write(
    
        "支持格式: TXT, CSV, PDF, Word, Excel, PPT, 图片(PDF, JPG, PNG等) | 使用 Deepseek API 和 PaddleOCR 提供技术支持")
    